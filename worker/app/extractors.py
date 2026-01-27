from __future__ import annotations

from typing import List, Tuple, Optional, Callable
import os

import pdfplumber
from docx import Document
from pptx import Presentation


def detect_mime_from_ext(file_path: str) -> Optional[str]:
    ext = os.path.splitext(file_path)[1].lower()
    return {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".txt": "text/plain",
        ".md": "text/markdown",
    }.get(ext)


def extract_pdf_pages(
    file_path: str,
    on_progress: Optional[Callable[[int, int], None]] = None
) -> Tuple[List[str], List[str]]:
    warnings: List[str] = []
    pages: List[str] = []
    with pdfplumber.open(file_path) as pdf:
        total = len(pdf.pages)
        for idx, p in enumerate(pdf.pages, start=1):
            pages.append(p.extract_text() or "")
            if on_progress:
                on_progress(idx, total)
    if len(pages) >= 3 and sum(len(p.strip()) for p in pages) < 800:
        warnings.append("low_text_density")
        warnings.append("possible_scanned_pdf")
    return pages, warnings


def extract_docx(file_path: str) -> Tuple[str, List[str]]:
    doc = Document(file_path)
    texts = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n".join(texts), []


def extract_pptx_slides(file_path: str) -> Tuple[List[str], List[str]]:
    prs = Presentation(file_path)
    slides: List[str] = []
    for slide in prs.slides:
        buf = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                buf.append(shape.text)
        slides.append("\n".join(buf))
    return slides, []


def extract_txt(file_path: str) -> Tuple[str, List[str]]:
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read(), []
    except UnicodeDecodeError:
        with open(file_path, "r", encoding="cp949", errors="ignore") as f:
            return f.read(), ["encoding_fallback_cp949"]
